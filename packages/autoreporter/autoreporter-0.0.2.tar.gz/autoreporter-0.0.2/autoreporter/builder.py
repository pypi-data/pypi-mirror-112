from io import StringIO
from PIL import Image
import matplotlib
import plotly

# https://stackoverflow.com/questions/18897511/how-to-drawimage-a-matplotlib-figure-in-a-reportlab-canvas

def fig2img(fig):
    imgdata = StringIO()
    fig.savefig(imgdata, format='png')
    imgdata.seek(0)  # rewind the data

    return Image.open(imgdata)

def builder(path, ref, fig_dict = dict(), var_dict = dict()):
    SCALE = 2
    PAD = 0.3

    from pptx import Presentation
    from pptx.util import Inches
    import tempfile
    import copy
    import os
    
    prs = Presentation(ref)
    out = Presentation()

    out.slide_width = prs.slide_width
    out.slide_height = prs.slide_height
    
    for slide in prs.slides:
        blank_slide_layout = out.slide_layouts[6]
        o_slide = out.slides.add_slide(blank_slide_layout)

        for shape in slide.shapes:
            try:
                if shape.text[0] == "!" and shape.text[1] == "(" and shape.text[-1] == ")":
                    f = tempfile.NamedTemporaryFile(suffix=".png")
                    fig = fig_dict[shape.text[2:-1]]
                    if type(fig) == matplotlib.figure.Figure:
                        fig.set_size_inches(SCALE*shape.width/Inches(1), SCALE*shape.height/Inches(1))
                        fig.set_tight_layout({"pad" : PAD})
                        fig.savefig(f.name, format='png',dpi = 300)
                        pic = o_slide.shapes.add_picture(f.name, shape.left, shape.top, width=shape.width)
                        continue
                    elif type(fig) == plotly.graph_objs._figure.Figure:
                        fig.write_image(f.name)
                        pic = o_slide.shapes.add_picture(f.name, shape.left, shape.top, width=shape.width)
                        continue

                fvar = False
                otxt = ""
                nvar = ""
                i = 0

                while i < len(shape.text):
                    if fvar:
                        if shape.text[i] != "]":
                            nvar += shape.text[i]
                        else:
                            otxt += var_dict[nvar]
                            fvar = False
                            nvar = ""
                    elif shape.text[i] == "$" and shape.text[i+1] == "[":
                        fvar = True
                        nvar = ""
                        i+=1 
                    else:
                        otxt += shape.text[i]
                    i+=1
                shape.text = otxt
                cp = copy.deepcopy(shape.element) 
                o_slide.shapes._spTree.insert_element_before(cp, 'p:extLst')
                continue
            except:
                pass
            
            try:
                img = shape.image.blob
                f = tempfile.NamedTemporaryFile()
                f.write(img)
                o_slide.shapes.add_picture(f.name, shape.left, shape.top, shape.width, shape.height)
                continue
            except:
                pass

            o_slide.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), 'p:extLst')
    out.save(path)
